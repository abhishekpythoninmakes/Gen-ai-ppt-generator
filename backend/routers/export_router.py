import json
import io
import logging
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from database import get_db
from models import User, Presentation
from auth import get_current_user

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/export", tags=["Export"])


def _hex_to_rgb(hex_color: str):
    """Convert hex color string to RGBColor."""
    from pptx.dml.color import RGBColor

    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return RGBColor(0x16, 0x16, 0x3F)  # fallback
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def _parse_color_with_alpha(color_str: str) -> tuple:
    """Return (RGBColor, alpha) where alpha is 0-1."""
    from pptx.dml.color import RGBColor
    if not color_str:
        return (RGBColor(0x00, 0x00, 0x00), 1.0)
    if color_str.startswith("rgba"):
        try:
            parts = color_str.replace("rgba", "").replace("(", "").replace(")", "").split(",")
            r = int(parts[0].strip())
            g = int(parts[1].strip())
            b = int(parts[2].strip())
            a = float(parts[3].strip())
            return (RGBColor(r, g, b), max(0.0, min(1.0, a)))
        except Exception:
            return (_hex_to_rgb("#000000"), 1.0)
    if color_str.startswith("#"):
        return (_hex_to_rgb(color_str), 1.0)
    return (_hex_to_rgb("#000000"), 1.0)


def _download_image(url: str) -> io.BytesIO | None:
    """Download an image from URL and return as BytesIO."""
    import httpx
    import base64
    import re

    if not url or url.startswith("data:"):
        if not url or not url.startswith("data:"):
            return None
        # Support data URLs (e.g., from uploads in the editor)
        try:
            match = re.match(r"^data:(.*?);base64,(.*)$", url, re.DOTALL)
            if not match:
                return None
            data_b64 = match.group(2)
            img_bytes = base64.b64decode(data_b64)
            buf = io.BytesIO(img_bytes)
            buf.seek(0)
            return buf
        except Exception as e:
            logger.warning(f"Failed to decode data URL image: {e}")
            return None
    try:
        with httpx.Client(timeout=10.0, follow_redirects=True) as client:
            resp = client.get(url)
            resp.raise_for_status()
            buf = io.BytesIO(resp.content)
            buf.seek(0)
            return buf
    except Exception as e:
        logger.warning(f"Failed to download image '{url[:80]}': {e}")
        return None


@router.get("/{ppt_id}/pptx")
def export_pptx(
    ppt_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """Export presentation as .pptx file."""
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    ppt = (
        db.query(Presentation)
        .filter(Presentation.id == ppt_id, Presentation.user_id == current_user.id)
        .first()
    )
    if not ppt:
        raise HTTPException(status_code=404, detail="Presentation not found")

    try:
        content = json.loads(ppt.content_json)
    except json.JSONDecodeError:
        raise HTTPException(status_code=400, detail="Invalid presentation data")

    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Scale factor: canvas 960x540 → 13.333"x7.5"
    SCALE_X = 13.333 / 960
    SCALE_Y = 7.5 / 540

    slides_data = content.get("slides", [])

    for i, slide_data in enumerate(slides_data):
        slide_layout = prs.slide_layouts[5]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set slide background color
        bg_color = slide_data.get("bg_color", "#1a1a2e")
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(bg_color)

        # Optional background image
        bg_image = slide_data.get("bg_image", "")
        if bg_image:
            img_buf = _download_image(bg_image)
            if img_buf:
                try:
                    slide.shapes.add_picture(
                        img_buf,
                        Inches(0),
                        Inches(0),
                        prs.slide_width,
                        prs.slide_height,
                    )
                except Exception as e:
                    logger.warning(f"Failed to add background image: {e}")

        # Check if slide has element-based data
        elements = slide_data.get("_elements", [])

        if elements:
            # New element-based rendering
            for el in elements:
                el_type = el.get("type", "")

                if el_type == "text":
                    x_in = Inches(el.get("x", 0) * SCALE_X)
                    y_in = Inches(el.get("y", 0) * SCALE_Y)
                    w_in = Inches(el.get("width", 300) * SCALE_X)
                    h_in = Inches(max(el.get("height", 50) * SCALE_Y, 0.3))

                    # Optional text background
                    bg_color = el.get("backgroundColor") or ""
                    bg_opacity = el.get("backgroundOpacity")
                    if bg_color:
                        rgb, alpha = _parse_color_with_alpha(bg_color)
                        if isinstance(bg_opacity, (int, float)):
                            alpha = max(0.0, min(1.0, float(bg_opacity)))
                        try:
                            from pptx.enum.shapes import MSO_SHAPE
                            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_in, y_in, w_in, h_in)
                            rect.fill.solid()
                            rect.fill.fore_color.rgb = rgb
                            rect.fill.transparency = max(0.0, min(1.0, 1.0 - alpha))
                            rect.line.fill.background()
                        except Exception as e:
                            logger.warning(f"Failed to add text background: {e}")

                    txBox = slide.shapes.add_textbox(x_in, y_in, w_in, h_in)
                    tf = txBox.text_frame
                    tf.word_wrap = True

                    text = el.get("text", "")
                    lines = text.split("\n")
                    for idx, line in enumerate(lines):
                        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                        # Strip bullet point markers
                        clean = line.lstrip("▸•- ").strip()
                        p.text = clean if clean else line

                        font_size = el.get("fontSize", 24)
                        p.font.size = Pt(font_size)

                        font_weight = el.get("fontWeight", "normal")
                        p.font.bold = font_weight == "bold"

                        font_style = el.get("fontStyle", "normal")
                        p.font.italic = font_style == "italic"

                        p.font.underline = el.get("underline", False)

                        # Color
                        fill_color = el.get("fill", "#ffffff")
                        if fill_color.startswith("rgba"):
                            # Parse rgba
                            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        else:
                            p.font.color.rgb = _hex_to_rgb(fill_color)

                        font_family = el.get("fontFamily", "Inter")
                        p.font.name = font_family

                        # Alignment
                        align = el.get("textAlign", "left")
                        if align == "center":
                            p.alignment = PP_ALIGN.CENTER
                        elif align == "right":
                            p.alignment = PP_ALIGN.RIGHT
                        else:
                            p.alignment = PP_ALIGN.LEFT

                        if idx > 0:
                            p.space_before = Pt(4)

                elif el_type == "image":
                    src = el.get("src", "")
                    if src:
                        img_buf = _download_image(src)
                        if img_buf:
                            x_in = Inches(el.get("x", 0) * SCALE_X)
                            y_in = Inches(el.get("y", 0) * SCALE_Y)
                            w_in = Inches(el.get("width", 300) * SCALE_X)
                            h_in = Inches(el.get("height", 200) * SCALE_Y)
                            try:
                                slide.shapes.add_picture(img_buf, x_in, y_in, w_in, h_in)
                            except Exception as e:
                                logger.warning(f"Failed to add image to slide: {e}")

                elif el_type == "shape":
                    x_in = Inches(el.get("x", 0) * SCALE_X)
                    y_in = Inches(el.get("y", 0) * SCALE_Y)
                    w_in = Inches(el.get("width", 200) * SCALE_X)
                    h_in = Inches(el.get("height", 200) * SCALE_Y)

                    from pptx.enum.shapes import MSO_SHAPE

                    shape_type = el.get("shapeType", "rect")
                    shape_map = {
                        "rect": MSO_SHAPE.ROUNDED_RECTANGLE,
                        "circle": MSO_SHAPE.OVAL,
                        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
                        "diamond": MSO_SHAPE.DIAMOND,
                        "arrow": MSO_SHAPE.RIGHT_ARROW,
                        "hexagon": MSO_SHAPE.HEXAGON,
                        "line": MSO_SHAPE.RECTANGLE,
                    }
                    mso = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)
                    shape = slide.shapes.add_shape(mso, x_in, y_in, w_in, h_in)

                    # Fill color
                    shape_fill = el.get("fill", "rgba(108, 99, 255, 0.3)")
                    if shape_fill.startswith("rgba"):
                        shape.fill.background()  # transparent
                    else:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = _hex_to_rgb(shape_fill)

                    # Stroke
                    stroke_color = el.get("stroke", "#6c63ff")
                    shape.line.color.rgb = _hex_to_rgb(stroke_color)
                    shape.line.width = Pt(el.get("strokeWidth", 2))

        else:
            # Legacy format fallback
            heading = slide_data.get("heading", "")
            if heading:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = heading
                p.font.size = Pt(36)
                p.font.bold = True
                text_color = slide_data.get("text_color", "#ffffff")
                p.font.color.rgb = _hex_to_rgb(text_color)
                p.alignment = PP_ALIGN.LEFT

            points = slide_data.get("points", [])
            if points:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.5), Inches(5))
                tf = txBox.text_frame
                tf.word_wrap = True
                for j, point in enumerate(points):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = f"• {point}"
                    p.font.size = Pt(20)
                    text_color = slide_data.get("text_color", "#ffffff")
                    p.font.color.rgb = _hex_to_rgb(text_color)
                    p.space_after = Pt(12)

            desc = slide_data.get("description", "")
            if desc:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(1))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = desc
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
                p.font.italic = True

            # Add image if present
            image_url = slide_data.get("image_url", "")
            if image_url:
                img_buf = _download_image(image_url)
                if img_buf:
                    try:
                        slide.shapes.add_picture(
                            img_buf, Inches(8.5), Inches(0.5), Inches(4.5), Inches(6)
                        )
                    except Exception as e:
                        logger.warning(f"Failed to add image to slide: {e}")

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    safe_title = "".join(c if c.isalnum() or c in " -_" else "" for c in ppt.title)[:50]
    filename = f"{safe_title or 'presentation'}.pptx"

    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
